#!/usr/bin/env python3
"""
Convert markdown presentation to PPTX with intelligent slide composition.

Design principles:
- One main idea per slide
- Max ~6 bullet points per slide
- Short, punchy text (not paragraphs)
- Clear visual hierarchy
- Tables get their own slides
"""

import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def parse_markdown(md_text: str) -> list[dict]:
    """Parse markdown into a list of slide specifications."""
    slides = []
    lines = md_text.split('\n')

    current_slide = None
    current_bullets = []
    in_table = False
    table_rows = []

    def flush_slide():
        nonlocal current_slide, current_bullets, in_table, table_rows
        if current_slide:
            if table_rows:
                current_slide['table'] = table_rows
                table_rows = []
            if current_bullets:
                current_slide['bullets'] = current_bullets
                current_bullets = []
            slides.append(current_slide)
            current_slide = None

    def flush_bullets_to_slide():
        nonlocal current_bullets
        if current_bullets and current_slide:
            if 'bullets' not in current_slide:
                current_slide['bullets'] = []
            current_slide['bullets'].extend(current_bullets)
            current_bullets = []

    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            i += 1
            continue

        # Skip horizontal rules
        if stripped == '---':
            i += 1
            continue

        # Level 1 header = title slide
        if stripped.startswith('# ') and not stripped.startswith('## '):
            flush_slide()
            title = stripped[2:].strip()
            current_slide = {'type': 'title', 'title': title, 'subtitle': ''}
            i += 1
            continue

        # Level 2 header = section slide
        if stripped.startswith('## '):
            flush_slide()
            title = stripped[3:].strip()
            # Clean up timing annotations like "(~25 min)"
            title = re.sub(r'\s*\(~?\d+\s*min\)', '', title)
            current_slide = {'type': 'section', 'title': title}
            i += 1
            continue

        # Level 3 header = content slide
        if stripped.startswith('### '):
            flush_slide()
            title = stripped[4:].strip()
            current_slide = {'type': 'content', 'title': title}
            i += 1
            continue

        # Bold line at start = could be a key point or subtitle
        if stripped.startswith('**') and stripped.endswith('**'):
            text = stripped[2:-2]
            if current_slide and current_slide['type'] == 'title':
                current_slide['subtitle'] = text
            elif current_slide:
                current_bullets.append({'text': text, 'bold': True, 'level': 0})
            i += 1
            continue

        # Table detection
        if '|' in stripped and stripped.startswith('|'):
            if not in_table:
                flush_bullets_to_slide()
                in_table = True
                table_rows = []

            # Skip separator rows
            if re.match(r'^\|[\s\-:|]+\|$', stripped):
                i += 1
                continue

            # Parse table row
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            if cells:
                table_rows.append(cells)
            i += 1
            continue
        else:
            if in_table:
                in_table = False
                if table_rows and current_slide:
                    current_slide['table'] = table_rows
                    table_rows = []

        # Bullet points
        if stripped.startswith('- ') or stripped.startswith('* '):
            text = stripped[2:].strip()
            level = 0
            # Check indentation for nested bullets
            indent = len(line) - len(line.lstrip())
            if indent >= 4:
                level = 1
            if indent >= 8:
                level = 2

            # Handle bold within bullet
            bold = False
            if text.startswith('**') and '**' in text[2:]:
                bold = True

            # Clean up markdown formatting
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Remove bold markers
            text = re.sub(r'\*([^*]+)\*', r'\1', text)  # Remove italic markers

            current_bullets.append({'text': text, 'bold': bold, 'level': level})
            i += 1
            continue

        # Numbered list
        if re.match(r'^\d+\.\s', stripped):
            text = re.sub(r'^\d+\.\s', '', stripped)
            current_bullets.append({'text': text, 'bold': False, 'level': 0, 'numbered': True})
            i += 1
            continue

        # Regular paragraph - treat as a bullet point if we're in a content slide
        if current_slide and current_slide['type'] == 'content':
            # Skip metadata lines
            if stripped.startswith('**') and ':' in stripped:
                # Key-value pair like **Date:** Monday...
                current_bullets.append({'text': stripped.replace('**', ''), 'bold': False, 'level': 0})
            elif not stripped.startswith('[') and not stripped.startswith('*['):
                # Skip links and footnotes, add regular text
                if len(stripped) > 10:  # Skip very short lines
                    current_bullets.append({'text': stripped, 'bold': False, 'level': 0})

        i += 1

    # Flush final slide
    flush_slide()

    return slides


def split_long_slides(slides: list[dict], max_bullets: int = 6) -> list[dict]:
    """Split slides with too many bullets into multiple slides."""
    result = []

    for slide in slides:
        if slide['type'] != 'content' or 'bullets' not in slide:
            result.append(slide)
            continue

        bullets = slide.get('bullets', [])
        if len(bullets) <= max_bullets:
            result.append(slide)
            continue

        # Split into multiple slides
        chunks = [bullets[i:i+max_bullets] for i in range(0, len(bullets), max_bullets)]
        for idx, chunk in enumerate(chunks):
            new_slide = slide.copy()
            new_slide['bullets'] = chunk
            if idx > 0:
                new_slide['title'] = f"{slide['title']} (cont.)"
            result.append(new_slide)

    return result


def create_presentation(slides: list[dict], output_path: str):
    """Create PPTX from slide specifications."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    for slide_spec in slides:
        if slide_spec['type'] == 'title':
            add_title_slide(prs, slide_spec)
        elif slide_spec['type'] == 'section':
            add_section_slide(prs, slide_spec)
        elif slide_spec['type'] == 'content':
            add_content_slide(prs, slide_spec)

    prs.save(output_path)
    print(f"Created {output_path} with {len(prs.slides)} slides")


def add_title_slide(prs: Presentation, spec: dict):
    """Add a title slide."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if spec.get('subtitle'):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = spec['subtitle']
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, spec: dict):
    """Add a section divider slide."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec['title']
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs: Presentation, spec: dict):
    """Add a content slide with title and bullets/table."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec['title']
    p.font.size = Pt(32)
    p.font.bold = True

    content_top = Inches(1.4)

    # Table
    if 'table' in spec and spec['table']:
        add_table(slide, spec['table'], content_top)
        return

    # Bullets
    if 'bullets' in spec and spec['bullets']:
        add_bullets(slide, spec['bullets'], content_top)


def add_bullets(slide, bullets: list[dict], top: float):
    """Add bullet points to slide."""
    content_box = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = bullet['text']
        p.font.size = Pt(22) if bullet.get('level', 0) == 0 else Pt(18)
        p.font.bold = bullet.get('bold', False)
        p.level = bullet.get('level', 0)
        p.space_before = Pt(8)
        p.space_after = Pt(4)


def add_table(slide, rows: list[list[str]], top: float):
    """Add a table to the slide."""
    if not rows:
        return

    num_cols = len(rows[0])
    num_rows = len(rows)

    # Calculate table dimensions
    table_width = Inches(11.5)
    table_height = Inches(0.4 * num_rows)
    left = Inches(0.9)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # Style the table
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            # Format
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(16)
            if row_idx == 0:
                para.font.bold = True
            para.alignment = PP_ALIGN.LEFT


def main():
    if len(sys.argv) < 2:
        print("Usage: uv run md2keynote.py <input.md> [output.pptx]")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        output_path = input_path.with_suffix('.pptx')

    md_text = input_path.read_text()
    slides = parse_markdown(md_text)
    slides = split_long_slides(slides, max_bullets=6)

    create_presentation(slides, str(output_path))


if __name__ == '__main__':
    main()
